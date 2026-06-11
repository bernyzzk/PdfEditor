from __future__ import annotations

from dataclasses import dataclass
import os
from pathlib import Path
import re
import tempfile
import winreg

import fitz
from PySide6.QtCore import QObject, Signal
from PySide6.QtGui import QImage


@dataclass
class PdfObject:
    kind: str
    rect: fitz.Rect
    text: str = ""
    font: str = "Helvetica"
    size: float = 11.0
    color: tuple[float, float, float] = (0.0, 0.0, 0.0)
    xref: int = 0
    font_xref: int = 0
    image_data: bytes = b""


class PdfDocument(QObject):
    changed = Signal()
    structure_changed = Signal()
    page_changed = Signal(int)
    dirty_changed = Signal(bool)

    def __init__(self) -> None:
        super().__init__()
        self.doc: fitz.Document | None = None
        self.path: Path | None = None
        self.page_index = 0
        self.dirty = False
        self._font_serial = 0
        self._undo: list[bytes] = []
        self._redo: list[bytes] = []
        self._render_cache: dict[tuple[int, float], QImage] = {}
        self._object_cache: dict[int, list[PdfObject]] = {}

    @property
    def is_open(self) -> bool:
        return self.doc is not None

    @property
    def page_count(self) -> int:
        return self.doc.page_count if self.doc else 0

    def open(self, path: str, password: str = "") -> bool:
        try:
            doc = fitz.open(path)
        except (OSError, RuntimeError, ValueError):
            return False
        if doc.needs_pass and not doc.authenticate(password):
            doc.close()
            return False
        for page in doc:
            if page.rotation:
                page.remove_rotation()
        self._replace_document(doc)
        self.path = Path(path)
        self.page_index = 0
        self._undo.clear()
        self._redo.clear()
        self._set_dirty(False)
        self.changed.emit()
        self.structure_changed.emit()
        self.page_changed.emit(0)
        return True

    def close(self) -> None:
        if self.doc:
            self.doc.close()
        self.doc = None
        self.path = None
        self.page_index = 0
        self._undo.clear()
        self._redo.clear()
        self._set_dirty(False)
        self.changed.emit()
        self.structure_changed.emit()

    def set_page(self, index: int) -> None:
        if not self.doc:
            return
        index = max(0, min(index, self.page_count - 1))
        if index != self.page_index:
            self.page_index = index
            self.page_changed.emit(index)

    def render_page(self, index: int, zoom: float = 1.0) -> QImage:
        if not self.doc:
            return QImage()
        key = (index, round(zoom, 3))
        cached = self._render_cache.get(key)
        if cached is not None:
            return cached
        page = self.doc.load_page(index)
        pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom), alpha=False)
        image = QImage(
            pix.samples,
            pix.width,
            pix.height,
            pix.stride,
            QImage.Format.Format_RGB888,
        )
        result = image.copy()
        self._render_cache[key] = result
        while len(self._render_cache) > 8:
            self._render_cache.pop(next(iter(self._render_cache)))
        return result

    def render_thumbnail(self, index: int, width: int = 130) -> QImage:
        if not self.doc:
            return QImage()
        page = self.doc.load_page(index)
        scale = width / page.rect.width
        return self.render_page(index, scale)

    def add_text(
        self,
        page_index: int,
        rect: fitz.Rect,
        text: str,
        font_size: float,
        color: tuple[float, float, float],
        font: str = "Helvetica",
    ) -> None:
        if not self.doc or not text:
            return
        self._checkpoint()
        page = self.doc.load_page(page_index)
        box = fitz.Rect(rect)
        fontname, fontfile = self._font_args(font)
        written = self._write_textbox(page, box, text, font_size, color, fontname, fontfile)
        if not written:
            page.insert_text(box.tl + (0, font_size), text, fontsize=font_size, fontname=fontname, fontfile=fontfile, color=color)
        self._commit()

    def page_objects(self, page_index: int) -> list[PdfObject]:
        if not self.doc:
            return []
        if page_index in self._object_cache:
            return self._object_cache[page_index]
        page = self.doc.load_page(page_index)
        objects: list[PdfObject] = []
        font_xrefs = {font[4]: font[0] for font in page.get_fonts(full=True)}
        image_xrefs: list[tuple[fitz.Rect, int]] = []
        try:
            image_xrefs = [
                (fitz.Rect(info["bbox"]), info.get("xref", 0))
                for info in page.get_image_info(xrefs=True)
                if info.get("xref", 0) and self.doc.xref_is_image(info["xref"])
            ]
        except Exception:
            pass
        data = page.get_text("dict")
        for block in data.get("blocks", []):
            if block.get("type") == 0:
                for line in block.get("lines", []):
                    spans = line.get("spans", [])
                    if not spans:
                        continue
                    rect = fitz.Rect(spans[0]["bbox"])
                    for span in spans[1:]:
                        rect |= fitz.Rect(span["bbox"])
                    first = spans[0]
                    objects.append(
                        PdfObject(
                            "text",
                            rect,
                            "".join(span.get("text", "") for span in spans),
                            self._clean_font(first.get("font", "Helvetica")),
                            float(first.get("size", 11)),
                            self._color_tuple(first.get("color", 0)),
                            font_xref=font_xrefs.get(first.get("font", ""), 0),
                        )
                    )
            elif block.get("type") == 1:
                if block.get("width", 0) <= 1 and block.get("height", 0) <= 1:
                    continue
                rect = fitz.Rect(block["bbox"])
                xref = next((xref for image_rect, xref in image_xrefs if image_rect.intersects(rect) and abs(image_rect.get_area() - rect.get_area()) < 2), 0)
                objects.append(PdfObject("image", rect, xref=xref, image_data=block.get("image", b"")))
        represented_images = [obj.rect for obj in objects if obj.kind == "image"]
        for rect, xref in image_xrefs:
            if not any(abs(existing.get_area() - rect.get_area()) < 2 and existing.intersects(rect) for existing in represented_images):
                objects.append(PdfObject("image", rect, xref=xref))
        try:
            page_area = page.rect.get_area()
            for drawing in page.get_drawings():
                rect = fitz.Rect(drawing["rect"])
                if rect.width >= 3 and rect.height >= 3 and 12 <= rect.get_area() < page_area * 0.9:
                    objects.append(PdfObject("graphic", rect))
        except Exception:
            pass
        self._object_cache[page_index] = objects
        return objects

    def object_at(self, page_index: int, point: fitz.Point) -> PdfObject | None:
        matches = [obj for obj in self.page_objects(page_index) if point in obj.rect]
        return min(matches, key=lambda obj: obj.rect.get_area()) if matches else None

    def replace_text(
        self,
        page_index: int,
        obj: PdfObject,
        text: str,
        font_size: float,
        color: tuple[float, float, float],
        font: str,
    ) -> None:
        if not self.doc:
            return
        self._checkpoint()
        page = self.doc.load_page(page_index)
        page.add_redact_annot(obj.rect + (-1, -1, 1, 1), fill=(1, 1, 1), cross_out=False)
        page.apply_redactions(images=0, graphics=0, text=0)
        self._compact_document()
        page = self.doc.load_page(page_index)
        if text:
            fontname, fontfile = self._font_for_object(page, obj, font)
            estimated_width = max(obj.rect.width, len(text) * font_size * 0.58)
            box = fitz.Rect(
                obj.rect.x0,
                obj.rect.y0 - 2,
                min(page.rect.x1 - 18, obj.rect.x0 + estimated_width + 8),
                obj.rect.y1 + max(6, font_size * 1.2),
            )
            self._write_textbox(page, box, text, font_size, color, fontname, fontfile)
        self._commit()

    def transform_object(self, page_index: int, obj: PdfObject, new_rect: fitz.Rect) -> None:
        if not self.doc or new_rect.is_empty:
            return
        self._checkpoint()
        page = self.doc.load_page(page_index)
        if obj.kind == "text":
            page.add_redact_annot(obj.rect + (-1, -1, 1, 1), fill=(1, 1, 1), cross_out=False)
            page.apply_redactions(images=0, graphics=0, text=0)
            self._compact_document()
            page = self.doc.load_page(page_index)
            fontname, fontfile = self._font_for_object(page, obj, obj.font)
            box = fitz.Rect(
                new_rect.x0,
                new_rect.y0 - 3,
                max(new_rect.x1, new_rect.x0 + len(obj.text) * obj.size * 0.65),
                max(new_rect.y1, new_rect.y0 + obj.size * 2.5),
            )
            self._write_textbox(page, box, obj.text, obj.size, obj.color, fontname, fontfile)
        elif obj.kind in ("image", "graphic"):
            # Render the visible composite so separate transparency masks and
            # unusual color spaces remain visually intact after moving.
            image_data = page.get_pixmap(
                matrix=fitz.Matrix(3, 3),
                clip=obj.rect,
                alpha=False,
            ).tobytes("png")
            page.add_redact_annot(obj.rect, fill=(1, 1, 1), cross_out=False)
            page.apply_redactions(images=2, graphics=2 if obj.kind == "graphic" else 0, text=1)
            self._compact_document()
            page = self.doc.load_page(page_index)
            page.insert_image(new_rect, stream=self._normalize_image(image_data), keep_proportion=False)
        self._commit()

    def copy_object_payload(self, page_index: int, obj: PdfObject) -> dict:
        payload = {
            "kind": obj.kind,
            "rect": tuple(obj.rect),
            "text": obj.text,
            "font": obj.font,
            "size": obj.size,
            "color": obj.color,
            "font_xref": obj.font_xref,
        }
        if self.doc and obj.kind != "text":
            page = self.doc.load_page(page_index)
            payload["image_data"] = page.get_pixmap(
                matrix=fitz.Matrix(3, 3), clip=obj.rect, alpha=False
            ).tobytes("png")
        return payload

    def paste_object(self, page_index: int, payload: dict, offset: float = 18.0) -> PdfObject | None:
        if not self.doc:
            return None
        source = fitz.Rect(payload["rect"])
        target = source + (offset, offset, offset, offset)
        page = self.doc.load_page(page_index)
        target = target & page.rect
        if target.is_empty:
            return None
        self._checkpoint()
        if payload["kind"] == "text":
            obj = PdfObject(
                "text",
                source,
                payload.get("text", ""),
                payload.get("font", "Helvetica"),
                float(payload.get("size", 11)),
                tuple(payload.get("color", (0, 0, 0))),
                font_xref=int(payload.get("font_xref", 0)),
            )
            fontname, fontfile = self._font_for_object(page, obj, obj.font)
            box = fitz.Rect(target.x0, target.y0 - 3, target.x1 + 8, target.y1 + obj.size * 1.5)
            self._write_textbox(page, box, obj.text, obj.size, obj.color, fontname, fontfile)
        else:
            image_data = payload.get("image_data", b"")
            if not image_data:
                self._undo.pop()
                return None
            page.insert_image(target, stream=self._normalize_image(image_data), keep_proportion=False)
        self._commit()
        detected = self.object_at(page_index, fitz.Point((target.x0 + target.x1) / 2, (target.y0 + target.y1) / 2))
        return detected or PdfObject(
            payload["kind"],
            target,
            payload.get("text", ""),
            payload.get("font", "Helvetica"),
            float(payload.get("size", 11)),
            tuple(payload.get("color", (0, 0, 0))),
            font_xref=int(payload.get("font_xref", 0)),
            image_data=payload.get("image_data", b""),
        )

    def erase_object(self, page_index: int, obj: PdfObject) -> None:
        if not self.doc:
            return
        self._checkpoint()
        page = self.doc.load_page(page_index)
        page.add_redact_annot(obj.rect + (-1, -1, 1, 1), fill=(1, 1, 1), cross_out=False)
        page.apply_redactions(images=2, graphics=1, text=0)
        self._commit()

    def add_image(self, page_index: int, rect: fitz.Rect, image_path: str) -> None:
        if not self.doc or rect.is_empty:
            return
        self._checkpoint()
        image_data = self._normalize_image(Path(image_path).read_bytes())
        self.doc.load_page(page_index).insert_image(rect, stream=image_data, keep_proportion=True)
        self._commit()

    def add_link(self, page_index: int, rect: fitz.Rect, uri: str) -> None:
        if not self.doc or rect.is_empty or not uri:
            return
        self._checkpoint()
        self.doc.load_page(page_index).insert_link({"kind": fitz.LINK_URI, "from": rect, "uri": uri})
        self._commit()

    def redact_text(self, query: str) -> int:
        if not self.doc or not query:
            return 0
        hits: list[tuple[int, fitz.Rect]] = []
        for index in range(self.page_count):
            for rect in self.doc.load_page(index).search_for(query):
                hits.append((index, rect))
        if not hits:
            return 0
        self._checkpoint()
        touched: set[int] = set()
        for index, rect in hits:
            self.doc.load_page(index).add_redact_annot(rect, fill=(0, 0, 0), cross_out=False)
            touched.add(index)
        for index in touched:
            self.doc.load_page(index).apply_redactions(images=0, graphics=0, text=0)
        self._commit()
        return len(hits)

    def replace_image(self, page_index: int, obj: PdfObject, image_path: str) -> None:
        if not self.doc or obj.kind != "image":
            return
        self._checkpoint()
        page = self.doc.load_page(page_index)
        page.add_redact_annot(obj.rect, fill=(1, 1, 1), cross_out=False)
        page.apply_redactions(images=2, graphics=0, text=1)
        page = self.doc.reload_page(page)
        image_data = self._normalize_image(Path(image_path).read_bytes())
        page.insert_image(obj.rect, stream=image_data, keep_proportion=True)
        self._commit()

    def add_ink(
        self,
        page_index: int,
        points: list[fitz.Point],
        color: tuple[float, float, float],
        width: float,
    ) -> None:
        if not self.doc or len(points) < 2:
            return
        self._checkpoint()
        stroke = [(point.x, point.y) for point in points]
        page = self.doc.load_page(page_index)
        annot = page.add_ink_annot([stroke])
        annot.set_colors(stroke=color)
        annot.set_border(width=width)
        annot.update()
        self._commit()

    def add_highlight(self, page_index: int, rect: fitz.Rect) -> None:
        if not self.doc or rect.is_empty or rect.width < 3 or rect.height < 3:
            return
        self._checkpoint()
        page = self.doc.load_page(page_index)
        annot = page.add_highlight_annot(rect)
        annot.set_colors(stroke=(1.0, 0.78, 0.12))
        annot.set_opacity(0.35)
        annot.update()
        self._commit()

    def add_signature(self, page_index: int, rect: fitz.Rect, image_path: str) -> None:
        self.add_image(page_index, rect, image_path)

    def add_stamp(self, page_index: int, rect: fitz.Rect, image_data: bytes) -> None:
        if not self.doc or rect.is_empty or not image_data:
            return
        self._checkpoint()
        self.doc.load_page(page_index).insert_image(rect, stream=image_data, keep_proportion=False, overlay=True)
        self._commit()

    def rotate_page(self, index: int, degrees: int = 90) -> None:
        if not self.doc:
            return
        self._checkpoint()
        page = self.doc.load_page(index)
        page.set_rotation((page.rotation + degrees) % 360)
        self._commit()

    def delete_page(self, index: int) -> None:
        if not self.doc or self.page_count <= 1:
            return
        self._checkpoint()
        self.doc.delete_page(index)
        self.page_index = min(index, self.page_count - 1)
        self._commit()
        self.structure_changed.emit()
        self.page_changed.emit(self.page_index)

    def move_page(self, source: int, destination: int) -> None:
        if not self.doc or source == destination:
            return
        self._checkpoint()
        order = list(range(self.page_count))
        moved = order.pop(source)
        order.insert(destination, moved)
        self.doc.select(order)
        self.page_index = destination
        self._commit()
        self.structure_changed.emit()
        self.page_changed.emit(destination)

    def merge_files(self, paths: list[str]) -> None:
        if not self.doc or not paths:
            return
        self._checkpoint()
        for path in paths:
            with fitz.open(path) as source:
                self.doc.insert_pdf(source)
        self._commit()
        self.structure_changed.emit()

    def extract_page(self, index: int, path: str) -> None:
        if not self.doc:
            return
        output = fitz.open()
        output.insert_pdf(self.doc, from_page=index, to_page=index)
        output.save(path, garbage=4, deflate=True)
        output.close()

    def run_ocr(self, language: str = "fra+eng", only_without_text: bool = True) -> int:
        if not self.doc:
            return 0
        tessdata = self._tessdata_path()
        if not tessdata:
            raise RuntimeError("Données Tesseract introuvables.")
        self._checkpoint()
        output = fitz.open()
        count = 0
        for index in range(self.page_count):
            page = self.doc.load_page(index)
            if only_without_text and page.get_text().strip():
                output.insert_pdf(self.doc, from_page=index, to_page=index)
                continue
            pix = page.get_pixmap(dpi=250, alpha=False)
            ocr_data = pix.pdfocr_tobytes(language=language, tessdata=tessdata)
            with fitz.open(stream=ocr_data, filetype="pdf") as ocr_page:
                output.insert_pdf(ocr_page)
            count += 1
        if count:
            self._replace_document(output)
            self._commit()
            self.structure_changed.emit()
        else:
            output.close()
            self._undo.pop()
        return count

    def export_word(self, path: str) -> None:
        if not self.doc:
            return
        from docx import Document
        from docx.shared import Pt

        document = Document()
        for index, page in enumerate(self.doc):
            if index:
                document.add_page_break()
            for obj in self.page_objects(index):
                if obj.kind != "text" or not obj.text.strip():
                    continue
                paragraph = document.add_paragraph()
                run = paragraph.add_run(obj.text)
                run.font.name = obj.font
                run.font.size = Pt(obj.size)
        document.save(path)

    def export_excel(self, path: str) -> None:
        if not self.doc:
            return
        from openpyxl import Workbook

        workbook = Workbook()
        workbook.remove(workbook.active)
        for index, page in enumerate(self.doc):
            sheet = workbook.create_sheet(f"Page {index + 1}")
            for row, line in enumerate(page.get_text("text").splitlines(), 1):
                sheet.cell(row, 1, line)
            sheet.column_dimensions["A"].width = 100
        workbook.save(path)

    def export_powerpoint(self, path: str) -> None:
        if not self.doc:
            return
        from io import BytesIO
        from pptx import Presentation
        from pptx.util import Inches

        presentation = Presentation()
        presentation.slide_width = Inches(11.69)
        presentation.slide_height = Inches(8.27)
        blank = presentation.slide_layouts[6]
        for page in self.doc:
            slide = presentation.slides.add_slide(blank)
            pix = page.get_pixmap(dpi=140, alpha=False)
            slide.shapes.add_picture(BytesIO(pix.tobytes("png")), 0, 0, presentation.slide_width, presentation.slide_height)
        presentation.save(path)

    def undo(self) -> None:
        if not self.doc or not self._undo:
            return
        self._redo.append(self._snapshot())
        self._restore(self._undo.pop())

    def redo(self) -> None:
        if not self.doc or not self._redo:
            return
        self._undo.append(self._snapshot())
        self._restore(self._redo.pop())

    def save(self, path: str) -> None:
        if not self.doc:
            return
        Path(path).write_bytes(self.doc.tobytes(garbage=4, deflate=True, clean=True))
        self.path = Path(path)
        self._set_dirty(False)

    def _checkpoint(self) -> None:
        self._undo.append(self._snapshot())
        if len(self._undo) > 20:
            self._undo.pop(0)
        self._redo.clear()

    def _snapshot(self) -> bytes:
        assert self.doc is not None
        return self.doc.tobytes(garbage=3, deflate=True)

    def _restore(self, data: bytes) -> None:
        page = self.page_index
        self._replace_document(fitz.open(stream=data, filetype="pdf"))
        self.page_index = min(page, self.page_count - 1)
        self._set_dirty(True)
        self.changed.emit()
        self.structure_changed.emit()
        self.page_changed.emit(self.page_index)

    def _replace_document(self, doc: fitz.Document) -> None:
        if self.doc:
            self.doc.close()
        self.doc = doc
        self._invalidate_caches()

    def _commit(self) -> None:
        self._invalidate_caches()
        self._set_dirty(True)
        self.changed.emit()

    def _invalidate_caches(self) -> None:
        self._render_cache.clear()
        self._object_cache.clear()

    def _set_dirty(self, value: bool) -> None:
        if value != self.dirty:
            self.dirty = value
            self.dirty_changed.emit(value)

    @staticmethod
    def _clean_font(font: str) -> str:
        return font.split("+", 1)[-1].replace("-", " ")

    @staticmethod
    def _color_tuple(value: int) -> tuple[float, float, float]:
        return (((value >> 16) & 255) / 255, ((value >> 8) & 255) / 255, (value & 255) / 255)

    @staticmethod
    def _font_args(font: str) -> tuple[str, str | None]:
        normalized = font.lower().replace(" ", "")
        aliases = {
            "helvetica": "helv", "arial": "helv", "timesnewroman": "tiro",
            "timesroman": "tiro", "couriernew": "cour", "courier": "cour",
        }
        if normalized in aliases:
            return aliases[normalized], None
        font_file = PdfDocument._windows_font_file(font)
        if font_file:
            resource_name = "F_" + re.sub(r"[^A-Za-z0-9]", "", font)[:24]
            return resource_name, font_file
        return "helv", None

    def _font_for_object(self, page: fitz.Page, obj: PdfObject, requested_font: str) -> tuple[str, str | None]:
        local_font = self._windows_font_file(requested_font)
        if local_font:
            self._font_serial += 1
            resource = "Edit_" + re.sub(r"[^A-Za-z0-9]", "", requested_font)[:14] + f"_{self._font_serial}"
            return resource, local_font

        if self.doc and obj.font_xref and self._clean_font(requested_font).lower() == self._clean_font(obj.font).lower():
            try:
                name, ext, _, content = self.doc.extract_font(obj.font_xref)
                if content:
                    base_resource = "Embedded_" + re.sub(r"[^A-Za-z0-9]", "", name)[:20]
                    self._font_serial += 1
                    resource = f"{base_resource}_{self._font_serial}"
                    cache = Path(tempfile.gettempdir()) / "pdfeditor-fonts"
                    cache.mkdir(exist_ok=True)
                    font_file = cache / f"{base_resource}-{obj.font_xref}.{ext or 'ttf'}"
                    if not font_file.exists():
                        font_file.write_bytes(content)
                    return resource, str(font_file)
            except (RuntimeError, ValueError):
                pass
        fontname, fontfile = self._font_args(requested_font)
        if fontfile:
            self._font_serial += 1
            unique = "Edit_" + re.sub(r"[^A-Za-z0-9]", "", requested_font)[:14] + f"_{self._font_serial}"
            return unique, fontfile
        fallback = self._windows_font_file("Arial")
        if fallback:
            self._font_serial += 1
            return f"Edit_Arial_{self._font_serial}", fallback
        return fontname, fontfile

    @staticmethod
    def _tessdata_path() -> str | None:
        import sys

        candidates = [
            Path(getattr(sys, "_MEIPASS", "")) / "ocr-data",
            Path(__file__).resolve().parent.parent / "ocr-data",
            Path(r"C:\Program Files\Tesseract-OCR\tessdata"),
        ]
        for path in candidates:
            if (path / "eng.traineddata").exists():
                return str(path)
        return None

    @staticmethod
    def _normalize_image(data: bytes) -> bytes:
        from io import BytesIO
        from PIL import Image

        source = BytesIO(data)
        output = BytesIO()
        with Image.open(source) as image:
            image.convert("RGB").save(output, format="PNG")
        return output.getvalue()

    @staticmethod
    def _normalize_text(text: str) -> str:
        return text.replace("\xa0", " ").replace("\u202f", " ")

    def _image_xref_at(self, page: fitz.Page, rect: fitz.Rect) -> int:
        if not self.doc:
            return 0
        try:
            for info in page.get_image_info(xrefs=True):
                xref = info.get("xref", 0)
                if xref and self.doc.xref_is_image(xref) and fitz.Rect(info["bbox"]).intersects(rect):
                    return xref
        except Exception:
            pass
        return 0

    def _compact_document(self) -> None:
        if not self.doc:
            return
        data = self.doc.tobytes(garbage=4, deflate=True, clean=True)
        self._replace_document(fitz.open(stream=data, filetype="pdf"))

    def _write_textbox(
        self,
        page: fitz.Page,
        rect: fitz.Rect,
        text: str,
        font_size: float,
        color: tuple[float, float, float],
        fontname: str,
        fontfile: str | None,
    ) -> bool:
        text = self._normalize_text(text)
        if fontfile:
            font = fitz.Font(fontfile=fontfile)
            writer = fitz.TextWriter(page.rect)
            overflow = writer.fill_textbox(
                rect,
                text,
                font=font,
                fontsize=font_size,
                lineheight=1.2,
            )
            writer.write_text(page, color=color, overlay=True)
            return not overflow
        return page.insert_textbox(
            rect,
            text,
            fontsize=font_size,
            fontname=fontname,
            color=color,
            lineheight=1.2,
            overlay=True,
        ) >= 0

    @staticmethod
    def _windows_font_file(font: str) -> str | None:
        wanted = re.sub(r"[^a-z0-9]", "", font.lower())
        for suffix in ("bolditalicmt", "boldmt", "italicmt", "regularmt", "mt"):
            if wanted.endswith(suffix):
                wanted = wanted[: -len(suffix)]
                break
        keys = (
            r"SOFTWARE\Microsoft\Windows NT\CurrentVersion\Fonts",
            r"SOFTWARE\WOW6432Node\Microsoft\Windows NT\CurrentVersion\Fonts",
        )
        for key_name in keys:
            try:
                with winreg.OpenKey(winreg.HKEY_LOCAL_MACHINE, key_name) as key:
                    for index in range(winreg.QueryInfoKey(key)[1]):
                        label, value, _ = winreg.EnumValue(key, index)
                        family = re.sub(r"[^a-z0-9]", "", label.lower().split("(")[0])
                        if wanted in family or family in wanted:
                            path = value if os.path.isabs(value) else os.path.join(os.environ.get("WINDIR", r"C:\Windows"), "Fonts", value)
                            if os.path.isfile(path):
                                return path
            except OSError:
                continue
        return None
